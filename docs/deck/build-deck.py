#!/usr/bin/env python3
"""Hiraia progress-update deck (early Sept 2026), 6 slides, brand-consistent (paper/ink/teal/gold, hand fonts).
Inputs: docs/deck/data/{coverage.json, draw-share-by-month.json, evidence.json}; screenshots in docs/deck/screenshots/.
Output: docs/deck/Hiraia-progress-update-2026-09.pptx (+ PDF via LibreOffice). Run with /tmp/deckenv/bin/python."""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
D = os.path.join(ROOT, "docs/deck"); DATA = os.path.join(D, "data"); SHOTS = os.path.join(D, "screenshots"); OUT = os.path.join(D, "out"); os.makedirs(OUT, exist_ok=True)
PAPER, INK, PRIMARY, ACCENT, COVER = "#fdfdf6", "#0c343d", "#165a6a", "#f2c14e", "#dde7e4"
MUTED = "#5f7a80"; CHART_TEAL, CHART_GOLD, CHART_GREY = "#1f7f93", "#c9961a", "#6b7a7d"; RAMP = ["#c5e1e6", "#8ac1cb", "#4a94a3", "#165a6a"]
F_TITLE, F_DISPLAY, F_BODY = "Mansalva", "Caveat Brush", "Patrick Hand"
cov = json.load(open(f"{DATA}/coverage.json")); share = json.load(open(f"{DATA}/draw-share-by-month.json"))
ev = json.load(open(f"{DATA}/evidence.json")) if os.path.exists(f"{DATA}/evidence.json") else {}
def rgb(h): return RGBColor.from_string(h.lstrip("#"))
# ---- matplotlib brand setup
for ttf in ("Mansalva-Regular.ttf", "CaveatBrush-Regular.ttf", "PatrickHand-Regular.ttf"):
    fp = os.path.expanduser(f"~/Library/Fonts/{ttf}")
    if os.path.exists(fp): font_manager.fontManager.addfont(fp)
plt.rcParams.update({"font.family": [F_BODY, "DejaVu Sans"], "font.size": 13, "axes.edgecolor": "#b9c6c9", "axes.labelcolor": INK, "xtick.color": INK, "ytick.color": INK, "axes.spines.top": False, "axes.spines.right": False, "figure.facecolor": PAPER, "axes.facecolor": PAPER, "savefig.facecolor": PAPER})
def chart_draw_share():
    import datetime as _dt
    d0 = _dt.date(2026, 6, 1); xs = [(_dt.date.fromisoformat(r["date"]) - d0).days for r in share]
    ticks = [(_dt.date(y, m, 1) - d0).days for (y, m) in [(2026, 6), (2026, 7), (2026, 8), (2026, 9), (2026, 10), (2026, 11), (2026, 12), (2027, 1), (2027, 2), (2027, 3), (2027, 4), (2027, 5)]]
    fig, ax = plt.subplots(figsize=(7.4, 3.9), dpi=200)
    og = [r["other grades"] for r in share]; ax.plot(xs, og, color=CHART_GREY, lw=1.6, ls="--", label="other grades' cells")
    ax.annotate("other grades' cells", (xs[2], og[2]), xytext=(0, 6), textcoords="offset points", ha="center", color=MUTED, fontsize=10)
    offs = {"G5 Q1": 7, "G5 Q2": 7, "G5 Q3": -13, "G5 Q4": 7}
    for i, q in enumerate(["G5 Q1", "G5 Q2", "G5 Q3", "G5 Q4"]):
        ys = [r[q] for r in share]; ax.plot(xs, ys, color=RAMP[i], lw=2.2, marker="o", ms=5, label=q)
        j = max(range(len(share)), key=lambda k: ys[k]); ax.annotate(q, (xs[j], ys[j]), xytext=(0, offs[q]), textcoords="offset points", ha="center", color=INK, fontsize=11)
    ax.set_xticks(ticks); ax.set_xticklabels(["Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May"]); ax.set_xlim(xs[0] - 10, xs[-1] + 10)
    ax.set_ylabel("share of draws (%)"); ax.set_ylim(0, 65); ax.grid(axis="y", color="#e3eaeb", lw=0.8)
    s0 = (_dt.date(2027, 4, 9) - d0).days; ax.axvspan(s0, xs[-1] + 10, color=COVER, alpha=0.6, lw=0); ax.text((s0 + xs[-1] + 10) / 2, 61, "summer", ha="center", fontsize=9.5, color=MUTED)
    ax.set_title("Grade-5 student: share of feed draws by cell, SY 2026–27 (untagged cards ≈1.5%, not shown)", fontsize=10.5, color=INK, loc="left", pad=8); ax.legend(frameon=False, fontsize=9, loc="upper left", bbox_to_anchor=(0.0, -0.16), ncol=5)
    p = f"{OUT}/draw-share.png"; fig.tight_layout(); fig.savefig(p); plt.close(fig); return p
def chart_cells():
    cells = cov["cells"]; grades = list(range(3, 11)); fig, ax = plt.subplots(figsize=(6.2, 3.6), dpi=200)
    w = 0.2
    for qi in range(4):
        vals = [cells.get(f"G{g}-Q{qi+1}", {}).get("cards", 0) for g in grades]
        ax.bar([g + (qi - 1.5) * w for g in grades], vals, width=w - 0.02, color=RAMP[qi], label=f"Q{qi+1}", linewidth=0)
    ax.set_xticks(grades); ax.set_xticklabels([f"G{g}" for g in grades]); ax.set_ylabel("illustrated cards in the pool"); ax.grid(axis="y", color="#e3eaeb", lw=0.8)
    ax.set_title("Cards per grade × curriculum quarter (any competency the card serves)", fontsize=12, color=INK, loc="left"); ax.legend(frameon=False, fontsize=10, ncol=4)
    p = f"{OUT}/cells.png"; fig.tight_layout(); fig.savefig(p); plt.close(fig); return p
def chart_runs(runs):
    runs = [r for r in runs if r.get("gpu_hours")]
    if not runs: return None
    runs = sorted(runs, key=lambda r: r["gpu_hours"])[-10:]
    fig, ax = plt.subplots(figsize=(6.0, 3.6), dpi=200)
    ax.barh([r["name"][:28] for r in runs], [r["gpu_hours"] for r in runs], color=CHART_TEAL, height=0.6, linewidth=0)
    for i, r in enumerate(runs): ax.text(r["gpu_hours"], i, f"  {r['gpu_hours']:.0f} h", va="center", fontsize=10, color=INK)
    ax.set_xlabel("GPU-hours (wall-clock × GPUs)"); ax.set_title("RunPod training and evaluation runs", fontsize=12, color=INK, loc="left"); ax.grid(axis="x", color="#e3eaeb", lw=0.8)
    p = f"{OUT}/runs.png"; fig.tight_layout(); fig.savefig(p); plt.close(fig); return p
def chart_metrics(metrics):
    pts = []
    for m in metrics:
        try: b, a = float(str(m["before"]).rstrip("%")), float(str(m["after"]).rstrip("%")); pts.append((m["metric"][:34], b, a))
        except Exception: pass
    if not pts: return None
    pts = pts[:6]; fig, ax = plt.subplots(figsize=(6.0, 3.6), dpi=200); y = range(len(pts)); h = 0.34
    ax.barh([i + h/2 for i in y], [p[1] for p in pts], height=h, color=CHART_GREY, label="before", linewidth=0)
    ax.barh([i - h/2 for i in y], [p[2] for p in pts], height=h, color=CHART_TEAL, label="after", linewidth=0)
    for i, p in enumerate(pts): ax.text(p[1], i + h/2, f" {p[1]:g}", va="center", fontsize=9, color=INK); ax.text(p[2], i - h/2, f" {p[2]:g}", va="center", fontsize=9, color=INK)
    ax.set_yticks(list(y)); ax.set_yticklabels([p[0] for p in pts]); ax.invert_yaxis(); ax.legend(frameon=False, fontsize=10); ax.grid(axis="x", color="#e3eaeb", lw=0.8)
    ax.set_title("Measured before → after (see sources in notes)", fontsize=12, color=INK, loc="left")
    p = f"{OUT}/metrics.png"; fig.tight_layout(); fig.savefig(p); plt.close(fig); return p
# ---- pptx helpers
COMMIT = "58d358495"
prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5); BLANK = prs.slide_layouts[6]
def new_slide(title, message, n):
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(PAPER)
    logo = os.path.join(ROOT, "packages/brand/logo.png")
    if os.path.exists(logo): s.shapes.add_picture(logo, Inches(0.45), Inches(0.32), height=Inches(0.55))
    text(s, 1.15, 0.28, 11.6, 0.7, title, F_TITLE, 30, PRIMARY)
    text(s, 0.5, 0.98, 12.3, 0.7, message, F_DISPLAY, 20, INK)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(12.33), Emu(9000)); line.fill.solid(); line.fill.fore_color.rgb = rgb(PRIMARY); line.line.fill.background()
    text(s, 0.5, 7.05, 9, 0.3, "Hiraia · progress update · early September 2026 · hiraia.org", F_BODY, 11, MUTED); text(s, 12.3, 7.05, 0.6, 0.3, str(n), F_BODY, 11, MUTED, align=PP_ALIGN.RIGHT)
    return s
def text(s, x, y, w, h, t, font, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(str(t).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align; r = p.add_run(); r.text = line; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return tb
def bullets(s, x, y, w, h, items, size=15, font=F_BODY):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); r = p.add_run(); r.text = "•  " + it; r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = rgb(INK); p.space_after = Pt(5)
    return tb
def stat(s, x, y, w, value, label):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.15)); box.fill.solid(); box.fill.fore_color.rgb = rgb(COVER); box.line.fill.background(); box.adjustments[0] = 0.12
    text(s, x, y + 0.05, w, 0.65, value, F_DISPLAY, 28, PRIMARY, align=PP_ALIGN.CENTER); text(s, x, y + 0.68, w, 0.42, label, F_BODY, 12, INK, align=PP_ALIGN.CENTER)
def shot(s, x, y, h, name, caption):
    p = os.path.join(SHOTS, name); w = h * 0.462
    if os.path.exists(p): s.shapes.add_picture(p, Inches(x), Inches(y), height=Inches(h))
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); ph.fill.solid(); ph.fill.fore_color.rgb = rgb(COVER); ph.line.color.rgb = rgb(PRIMARY); ph.line.width = Pt(1); ph.adjustments[0] = 0.08
        text(s, x + 0.08, y + h/2 - 0.85, w - 0.16, 1.7, f"placeholder — screenshot of the unpublished question-cards build (commit {COMMIT}); not in the published v0.1 APK", F_BODY, 9.5, MUTED, align=PP_ALIGN.CENTER)
    text(s, x - 0.3, y + h + 0.05, w + 0.6, 0.4, caption, F_BODY, 11, MUTED, align=PP_ALIGN.CENTER)
def table(s, x, y, w, rows, col_w, size=11):
    n, m = len(rows), len(rows[0]); t = s.shapes.add_table(n, m, Inches(x), Inches(y), Inches(w), Inches(0.28 * n)).table
    for j, cw in enumerate(col_w): t.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.text = str(val); c.fill.solid(); c.fill.fore_color.rgb = rgb(COVER if i == 0 else PAPER); c.margin_left = c.margin_right = Inches(0.06); c.margin_top = c.margin_bottom = Inches(0.02)
            for p in c.text_frame.paragraphs:
                for r in p.runs: r.font.name = F_BODY; r.font.size = Pt(size); r.font.bold = i == 0; r.font.color.rgb = rgb(INK)
    return t
def notes(s, txt): s.notes_slide.notes_text_frame.text = txt
# ---- slides
brand = ev.get("brand", {}); runs = ev.get("runs", {}); evals = ev.get("evals", {})
import re as _re
import re
def _url(v, default): m = _re.search(r"https?://[^\s)\]'\"]+", str(v or "")); return m.group(0).rstrip(".,;") if m else default
_L = brand.get("links", {})
links = { "website": _url(_L.get("website"), "https://hiraia.org"), "apk_download": _url(_L.get("apk_download"), "https://hiraia.org/models/hiraia.apk"), "github": _url(_L.get("github"), "https://github.com/helloluis/hiraia"), "huggingface": _url(_L.get("huggingface"), "") }
hf_public = "NONE PUBLIC" not in str(_L.get("huggingface", "")).upper() and bool(links["huggingface"])
_a = str(brand.get("latest_apk", "")); m = _re.search(r"version\s*'?([\d.]+)'?", _a); d = _re.search(r"(\d{1,2} \w{3} 20\d\d)", _a)
apk = f"app v{m.group(1)}" + (f", published {d.group(1)}" if d else "") if m else "current build"
# 1 title
s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(PAPER)
logo = os.path.join(ROOT, "packages/brand/logo.png")
if os.path.exists(logo): s.shapes.add_picture(logo, Inches(0.9), Inches(0.8), height=Inches(1.1))
text(s, 0.9, 2.1, 11.5, 1.2, "Hiraia — progress update", F_TITLE, 54, PRIMARY); text(s, 0.9, 3.2, 11.5, 0.7, "Early September 2026", F_DISPLAY, 30, INK)
text(s, 0.9, 4.0, 11.5, 1.2, "An offline, on-device science tutor for Filipino grade-school students (Grades 3–10), aligned to the DepEd MATATAG curriculum, in Tagalog, Cebuano and English. The published build (v0.1) targets Android 12+ phones with 6 GB+ memory; a 1B budget tier runs on a Redmi (Snapdragon 685) in unpublished builds.", F_BODY, 17, INK)
for i, (v, l) in enumerate([(f"{cov['pool_cards']:,}", "illustrated fact cards"), (f"{cov['summary']['bank_facts']:,}", "verified trilingual facts"), (f"{cov['summary']['competencies']}", f"MATATAG competencies tracked ({cov.get('competencies_with_cards', 0)} with cards)"), (f"{cov['pool_cards_with_mcq']:,}", "cards with a quiz question")]): stat(s, 0.9 + i * 3.05, 5.35, 2.85, v, l)
notes(s, f"Sources: rag/bank/COVERAGE-ROUNDUP.md/.json (pool cards, bank facts, competencies with cards); packages/mobile/src/data/cards-questions.json (cards with an MCQ = {cov['pool_cards_with_mcq']:,}); repository commit {COMMIT} on branch question-cards; published APK: hiraia.org/models/hiraia.apk = v0.1.0 (June 2026, Sailor2-3B tier).")
text(s, 0.9, 6.7, 11.5, 0.4, f"Evidence: repository at commit {COMMIT} (branch question-cards) and its evaluation records; the published app is {apk} — the feed shown here ships in the next build. Sources are in the slide notes.", F_BODY, 11, MUTED)
# 2 runs
rc = ev.get("runs_curated", {}); RC = rc.get("rows", [])
s = new_slide("Training runs on RunPod", "Most of the compute went into a continued-pretrained Filipino base model; the rest into tutoring fine-tunes, a pruning experiment, and evaluation.", 2)
rows = [["run", "GPUs", "GPU-h", "cost", "basis", "outcome"]] + [[r["name"], r["gpu"], f"{r['hours']:.0f}", f"${r['cost']:,.0f}", r["basis"], r["outcome"]] for r in RC]
if len(rows) == 1: rows.append(["(pending)", "", "", "", "", ""])
table(s, 0.5, 1.9, 12.3, rows, [3.5, 1.4, 0.7, 0.8, 1.6, 4.3], size=11)
text(s, 0.5, 6.05, 12.3, 0.95, rc.get("totals", ""), F_DISPLAY, 14, PRIMARY)
_cav = " ".join(x for x in re.split(r"(?<=[.;]) ", runs.get("caveats", "")) if "memory" not in x.lower())
notes(s, "Per-run sources and hour/cost bases: docs/deck/data/evidence.json (runs, 49 entries; per-entry source = repo file:line or git commit). " + rc.get("grouping", "") + " " + _cav)
# 3 feed
s = new_slide("A curriculum-aware feed", "Cards are weighted by the student's grade and the curriculum quarter inferred from the device date; the boost moves through the school year by itself.", 3)
s.shapes.add_picture(chart_draw_share(), Inches(0.5), Inches(1.9), width=Inches(7.4))
shot(s, 8.25, 1.9, 3.2, "02-feed-card.png", "fact card"); shot(s, 10.35, 1.9, 3.2, "04-quiz-card.png", "interject quiz card")
w = cov["weights"]; table(s, 0.5, 5.65, 7.4, [["current quarter, same grade", "adjacent quarter", "other quarter, same grade", "adjacent grade, current", "off-curriculum"], [f"×{w['current']}", f"×{w['adjacent']}", f"×{w['other_same_grade']}", f"×{w['adjacent_grade_current']}", f"×{w['off']}"]], [1.6, 1.3, 1.6, 1.5, 1.4], size=10)
bullets(s, 8.25, 5.45, 4.8, 1.5, [f"{cov['pool_cards']:,} illustrated cards; {cov['summary']['bank_facts_labelled']:,} facts labelled to {cov['summary']['competencies']} competencies (LLM labels, 83% inter-model agreement; single-model cells down-weighted)", f"{cov['pool_cards_with_mcq']:,} cards ({100*cov['pool_cards_with_mcq']//cov['pool_cards']}%) carry a verified quiz question", f"0 elementary competencies below floor; {cov['summary']['below_floor_now']} junior-high competencies below floor in the card pool ({len(cov['laggards'])} once DepEd module cards are counted)"], size=12)
notes(s, f"Weighting: rag/pipeline/FEED-WEIGHTING.md; simulation from packages/shared/src/curriculum/feedWeighting.ts on the bundled tags (docs/deck/data/draw-share-by-month.json; the chart attributes a multi-cell card to the Grade-5 cell that won, so its Aug figure (26.1%) is not the same statistic as feed-calibration.mts's 'draws carrying a G5-Q2 cell' ({cov['calibration_g5_q2_share']}%) — both come from the same module and pool). Quarter weighting applies to all 324 competencies (elementary and junior high). Coverage: rag/bank/COVERAGE-ROUNDUP.md (below floor now = {cov['summary']['below_floor_now']}, all junior high; laggards after module cards = {len(cov['laggards'])}). Quiz: packages/mobile/src/data/cards-questions.json ({cov['pool_cards_with_mcq']:,} of {cov['pool_cards']:,}).")
# 4 model
ec = ev.get("evals_curated", {})
s = new_slide("Model progress: fluency and pedagogy", ec.get("message", ""), 4)
if ec.get("ppl"):
    fig, ax = plt.subplots(figsize=(5.6, 3.5), dpi=200); langs = [x[0] for x in ec["ppl"]]; y = range(len(langs)); h = 0.34
    ax.barh([i + h/2 for i in y], [x[1] for x in ec["ppl"]], height=h, color=CHART_GREY, label="base model", linewidth=0)
    ax.barh([i - h/2 for i in y], [x[2] for x in ec["ppl"]], height=h, color=CHART_TEAL, label="after continued pre-training", linewidth=0)
    for i, x in enumerate(ec["ppl"]): ax.text(x[1], i + h/2, f" {x[1]:.1f}", va="center", fontsize=10, color=INK); ax.text(x[2], i - h/2, f" {x[2]:.1f}", va="center", fontsize=10, color=INK)
    ax.set_yticks(list(y)); ax.set_yticklabels(langs); ax.invert_yaxis(); ax.set_xlabel("held-out perplexity (lower is better)"); ax.legend(frameon=False, fontsize=9.5, loc="lower right"); ax.set_xlim(0, 38); ax.grid(axis="x", color="#e3eaeb", lw=0.8)
    ax.set_title("Qwen3.5-2B before and after continued pre-training on Filipino text\n" + ec.get("ppl_caveat", ""), fontsize=9.5, color=INK, loc="left", pad=8)
    cp = f"{OUT}/ppl.png"; fig.tight_layout(); fig.savefig(cp); plt.close(fig); s.shapes.add_picture(cp, Inches(0.5), Inches(1.9), width=Inches(5.3))
text(s, 0.5, 5.25, 5.7, 0.3, "Tiles: pedagogy figures are from the June 3B line; language accuracy and hedging from the August 2B (SFT v1). The 2B capability run is pending.", F_BODY, 9.5, MUTED)
tiles = ec.get("tiles", [])
for i, (v, l) in enumerate(tiles):
    x = 0.5 + i * 1.92; box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.6), Inches(1.8), Inches(1.32)); box.fill.solid(); box.fill.fore_color.rgb = rgb(COVER); box.line.fill.background(); box.adjustments[0] = 0.12
    text(s, x, 5.63, 1.8, 0.5, v, F_DISPLAY, 20, PRIMARY, align=PP_ALIGN.CENTER); text(s, x + 0.05, 6.1, 1.7, 0.8, l, F_BODY, 9, INK, align=PP_ALIGN.CENTER)
text(s, 6.4, 1.9, 5.0, 0.4, "Pedagogy (measured on scripted probes)", F_DISPLAY, 15, PRIMARY); text(s, 6.4, 2.3, 5.0, 1.9, ec.get("pedagogy", ""), F_BODY, 12.5, INK)
text(s, 6.4, 4.25, 5.0, 0.4, "Known weaknesses (measured)", F_DISPLAY, 15, PRIMARY); bullets(s, 6.4, 4.65, 5.0, 2.2, ec.get("weaknesses", []), size=11.5)
shot(s, 11.65, 1.9, 3.0, "05-settings.png", "grade + language settings")
notes(s, "Sources: " + ec.get("sources", "") + "\n" + "\n".join(f"{m['metric']}: {m['before']} → {m['after']} — {m['source']}" for m in evals.get("metrics", [])))
# 5 next
s = new_slide("Next 2–3 weeks", "Consolidate the model, then prove it on the device.", 5)
left = ["More SFT and knowledge distillation on the CPT base (mode-locked buckets, teacher-generated tutoring turns)", "Extensive synthetic testing on-device: scripted student sessions on the Redmi, TTFT and per-page latency budgets", "Model-as-judge quality control on every generated fact, quiz and translation (decorrelated judge, AUP-safe routing)", "APK-level optimisations: bundle size (engravings subset, indexed PNG), cold start, SQLite migrations, edge-walk cost"]
right = ["Native-speaker review of the Tagalog and Cebuano copy (onboarding, settings, hedging clause)", "Junior-high gap fill: Lane B facts for the 16 remaining thin competencies; module-card supply", "Finish the quiz lane (10.7k pool cards still without a question) and per-competency Miss-Card labels", "Over-the-air fact-bank updates so content ships without an APK", "Retrieval-QA cases per new competency before each gate run", "Draft a classroom-pilot protocol (consent, offline logging, pre/post measures) to propose to DepEd schools"]
bullets(s, 0.5, 1.95, 6.3, 5, left, size=16); bullets(s, 6.9, 1.95, 6.0, 5, right, size=16)
text(s, 0.5, 6.1, 12.3, 0.4, "To report next time: the 2B capability run on the same 143 probes; on-device TTFT and per-page latency on the Redmi; quiz coverage after the lane completes; bundle size after optimisation.", F_DISPLAY, 14, PRIMARY)
notes(s, "Plan items trace to: rag/pipeline/FACT-SWARM-SPEC.md (Lane B, quiz lane, module cards), rag/pipeline/FEED-WEIGHTING.md (Miss-Card labels, edge-walk cost), finetuning/eval/routing/README.md (hedge-clause wording sweep), finetuning/cpt/*.md (SFT/KD on the CPT base), memory of on-device measurements on the Redmi SD685 (kitten tier ~4.4 tok/s, TTFT ~24 s).")
# 6 links
s = new_slide("Links", "Code, build and evaluation records are public; model weights follow once the release checklist is done.", 6)
L = [("Website", links["website"], ""), ("APK (direct download)", links["apk_download"], f"v0.1.0, 325 MB, sha256 708ea605…, Android 12+, 6 GB+ RAM; mirror: github.com/helloluis/hiraia/releases/tag/v0.1.0"), ("GitHub", links["github"], "source, pipelines and evaluation harness (public)"),
     ("Hugging Face", links["huggingface"] if hf_public else "to be published", "model weights + evaluation kit: release pending; corpus provenance will be documented, DepEd-derived text is not redistributed" if not hf_public else "")]
for i, (k, v, note) in enumerate(L):
    text(s, 0.8, 2.0 + i * 1.15, 3.2, 0.6, k, F_DISPLAY, 22, PRIMARY); tb = text(s, 4.0, 2.05 + i * 1.15, 8.6, 0.6, v, F_BODY, 20, INK if v.startswith("http") else MUTED)
    if v.startswith("http"):
        _r = tb.text_frame.paragraphs[0].runs[0]; _r.hyperlink.address = v; _r.font.color.rgb = rgb("#0f4a56")
    if note: text(s, 4.0, 2.5 + i * 1.15, 8.6, 0.4, note, F_BODY, 12, MUTED)
notes(s, f"Website and APK verified live on 2026-08-28 (APK v0.1.0, 324.8 MB, published June 2026; Android 12+, 6 GB+ recommended). GitHub repo is public. Hugging Face: model repos under the project account are private pending a release checklist (weights + evaluation kit); DepEd-derived corpus text will not be redistributed.")
p = f"{D}/Hiraia-progress-update-2026-09.pptx"; prs.save(p); print("wrote", p)
